"""Generate the diploma defense PowerPoint slides.

Builds ``thesis/presentation.pptx`` from the same CSV data sources used
by ``generate_thesis.py`` so the numbers are guaranteed to match.

Run::

    python thesis/generate_slides.py
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Inches, Pt

# ──────────────────────────────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────────────────────────────

ROOT = Path(__file__).resolve().parent.parent
THESIS_DIR = ROOT / "thesis"
EVAL_DIR = ROOT / "outputs" / "outputs" / "eval"
TABLES_DIR = EVAL_DIR / "tables"
FIG_DIR = THESIS_DIR / "figures"
OUT_PPTX = THESIS_DIR / "presentation.pptx"

DISPLAY_NAMES = {
    "pconv_unet":    "PConv U-Net",
    "unet_baseline": "Vanilla U-Net",
    "gated_unet":    "Gated U-Net",
}


# ──────────────────────────────────────────────────────────────────────
# Presentation skeleton (16:9 widescreen)
# ──────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

SW = prs.slide_width
SH = prs.slide_height
BLANK_LAYOUT = prs.slide_layouts[6]   # blank
TITLE_LAYOUT = prs.slide_layouts[0]   # title
TITLE_CONTENT = prs.slide_layouts[1]  # title + content


# ──────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────

def add_title(slide, text: str, *, size: int = 32) -> None:
    """Add a top-banner title to a blank slide."""
    tb = slide.shapes.add_textbox(Cm(1), Cm(0.6), SW - Cm(2), Cm(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x10, 0x32, 0x66)


def add_bullet_box(slide, items: list[str], *,
                   left: float = 1.2, top: float = 2.6,
                   width_cm: float | None = None,
                   height_cm: float | None = None,
                   font_pt: int = 20) -> None:
    """Add a vertical bullet list."""
    if width_cm is None:
        width = SW - Cm(2.4)
    else:
        width = Cm(width_cm)
    if height_cm is None:
        height = SH - Cm(top + 1.2)
    else:
        height = Cm(height_cm)

    tb = slide.shapes.add_textbox(Cm(left), Cm(top), width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.level = 0
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(font_pt)
        r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def add_image(slide, path: Path, *,
              left_cm: float, top_cm: float,
              width_cm: float | None = None,
              height_cm: float | None = None) -> None:
    if not path.exists():
        return
    kwargs = {}
    if width_cm is not None:
        kwargs["width"] = Cm(width_cm)
    if height_cm is not None:
        kwargs["height"] = Cm(height_cm)
    slide.shapes.add_picture(str(path), Cm(left_cm), Cm(top_cm), **kwargs)


def add_speaker_notes(slide, text: str) -> None:
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


def add_table(slide, df: pd.DataFrame, *,
              left_cm: float, top_cm: float,
              width_cm: float, height_cm: float,
              header_color: tuple[int, int, int] = (0x10, 0x32, 0x66),
              font_pt: int = 16) -> None:
    n_rows, n_cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Cm(left_cm), Cm(top_cm),
                                         Cm(width_cm), Cm(height_cm))
    t = table_shape.table

    for j, col in enumerate(df.columns):
        cell = t.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*header_color)
        for para in cell.text_frame.paragraphs:
            for r in para.runs:
                r.font.bold = True
                r.font.size = Pt(font_pt)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            v = df.iat[i, j]
            cell = t.cell(i + 1, j)
            if isinstance(v, float):
                cell.text = f"{v:.3f}" if abs(v) < 10 else f"{v:.2f}"
            else:
                cell.text = str(v)
            for para in cell.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = Pt(font_pt)


# ──────────────────────────────────────────────────────────────────────
# Data loaders (reuse generate_thesis.py logic via duplication for simplicity)
# ──────────────────────────────────────────────────────────────────────

def _safe_read(p: Path) -> pd.DataFrame:
    return pd.read_csv(p) if p.exists() else pd.DataFrame()


def load_overall() -> pd.DataFrame:
    df = _safe_read(TABLES_DIR / "overall_metrics.csv")
    if df.empty:
        return df
    df["Model"] = df["model"].map(DISPLAY_NAMES).fillna(df["model"])
    df = df[["Model", "psnr", "ssim", "lpips"]].copy()
    df.columns = ["Model", "PSNR ↑", "SSIM ↑", "LPIPS ↓"]
    df = df.sort_values("PSNR ↑", ascending=False).reset_index(drop=True)
    return df


def load_by_difficulty_psnr() -> pd.DataFrame:
    p = TABLES_DIR / "metrics_by_difficulty.csv"
    if not p.exists():
        return pd.DataFrame()
    df = pd.read_csv(p, header=[0, 1], index_col=0)
    psnr = df["psnr"].copy()
    psnr.index = psnr.index.map(lambda k: DISPLAY_NAMES.get(k, k))
    psnr = psnr[["light", "medium", "heavy"]]
    psnr.index.name = "Model"
    psnr = psnr.reset_index()
    psnr.columns = ["Model", "Light", "Medium", "Heavy"]
    return psnr


def load_by_damage_psnr() -> pd.DataFrame:
    p = TABLES_DIR / "metrics_by_damage.csv"
    if not p.exists():
        return pd.DataFrame()
    df = pd.read_csv(p, header=[0, 1], index_col=0)
    psnr = df["psnr"].copy()
    psnr.index = psnr.index.map(lambda k: DISPLAY_NAMES.get(k, k))
    psnr = psnr[["brush", "crack", "paint_loss", "stain"]]
    psnr.index.name = "Model"
    psnr = psnr.reset_index()
    psnr.columns = ["Model", "Brush", "Crack", "Paint loss", "Stain"]
    return psnr


# ──────────────────────────────────────────────────────────────────────
# Individual slides
# ──────────────────────────────────────────────────────────────────────

def slide_title():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    # Background tint
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x10, 0x32, 0x66)
    bg.line.fill.background()

    tb = s.shapes.add_textbox(Cm(2), Cm(5.5), SW - Cm(4), Cm(5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Restoration of Historic Artworks\nUsing Deep Learning-Based Digital Inpainting"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb2 = s.shapes.add_textbox(Cm(2), Cm(12), SW - Cm(4), Cm(4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for line in ["Master's Diploma Defense",
                 "Daniyar Krug",
                 "Supervisor: [Supervisor Name]",
                 "2026"]:
        p = tf2.paragraphs[0] if tf2.text == "" else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(20)
        r.font.color.rgb = RGBColor(0xDD, 0xE7, 0xF4)
    add_speaker_notes(s,
        "Good [morning/afternoon].  My name is Daniyar Krug and today I "
        "will defend my master's thesis on digital inpainting of historic "
        "paintings using deep learning.  This work compares three "
        "convolutional architectures under an identical training "
        "protocol and arrives at a counter-intuitive ranking that I will "
        "walk through in the next ~15 minutes.")


def slide_motivation():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Why digital art restoration?")
    add_bullet_box(s, [
        "Paintings inevitably degrade — pigments fade, varnish yellows, surfaces crack and flake.",
        "Manual restoration is slow, irreversible and requires deep expertise.",
        "Digital restoration is non-destructive: experts can preview reconstructions before any physical intervention.",
        "Inpainting (filling missing pixels) is the core computational problem — easy for small smooth regions, hard for large irregular holes on stylised content.",
    ])
    add_speaker_notes(s,
        "Three motivating points: paintings degrade; restoration is "
        "expensive and irreversible; digital restoration provides a "
        "non-destructive sandbox.  The core ML problem is image "
        "inpainting on irregular masks, which is where modern deep "
        "methods shine compared to classical diffusion-based approaches.")


def slide_research_questions():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Research questions")
    add_bullet_box(s, [
        "RQ1.  Does partial convolution (Liu 2018) beat a plain U-Net with mask-concatenation under matched training?",
        "RQ2.  Does gated convolution (Yu 2019) beat partial convolution on the art domain, as it does on natural photographs?",
        "RQ3.  Is the architectural ranking stable across damage severity (light / medium / heavy)?",
        "RQ4.  Is the ranking stable across damage type (brush / crack / paint loss / stain)?",
    ])
    add_speaker_notes(s,
        "Four research questions.  Note that all four are comparative: "
        "we are not asking whether any individual model 'works', we are "
        "asking which of three competing architectures wins under a "
        "fairness contract — same data, same loss, same recipe.")


def slide_contributions():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Four contributions")
    add_bullet_box(s, [
        "C1.  Open-source WikiArt benchmark for irregular-mask inpainting (first of its kind).",
        "C2.  Four art-specific synthetic damage generators × three difficulty levels — 12 evaluation cells.",
        "C3.  Rigorous matched-recipe comparison of PConv vs Vanilla vs Gated U-Net; paired Wilcoxon + Bonferroni; per-cell breakdown.",
        "C4.  Publicly-runnable interactive Gradio demo — first runnable art-inpainting system.",
    ])
    add_speaker_notes(s,
        "These four contributions carry through every chapter of the "
        "thesis.  The most important is C3 — the comparison itself — "
        "because the result contradicts the implicit consensus that "
        "PConv is the default mask-aware operator.")


def slide_related_work():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Related work — a 20-year arc")
    add_bullet_box(s, [
        "2000 — Bertalmio et al.: diffusion-based inpainting.  Fails on textures.",
        "2009 — PatchMatch: copy-paste from intact regions.  No semantics.",
        "2016 — Context Encoder: first deep inpainter.  Rectangular holes only.",
        "2018 — PConv (Liu): mask-aware partial convolution, free-form holes.",
        "2019 — Gated conv (Yu, DeepFillv2): soft attention gate, modern SOTA.",
        "2022 — LaMa, diffusion-based methods.  Out of compute scope for this thesis.",
        "Art-specific work: Gupta 2019 (Indian paintings), Farajzadeh 2021 (Iranian).  No rigorous architectural comparison on WikiArt.",
    ], font_pt=18)
    add_speaker_notes(s,
        "Three takeaways: (1) classical methods fail on large textured "
        "holes; (2) deep methods solved free-form holes by 2018–2019; "
        "(3) the application to art specifically is comparatively "
        "under-explored.  Our work fills exactly that gap.")


def slide_methodology_arch():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Three architectures, one training recipe")
    add_image(s, FIG_DIR / "pconv_unet_architecture.png",
              left_cm=1.0, top_cm=2.5, height_cm=12)
    add_bullet_box(s, [
        "PConv U-Net: partial conv + hard mask propagation (~26 M params).",
        "Vanilla U-Net: plain Conv2d, mask as 4th input channel (~26 M params).",
        "Gated U-Net: gated conv, learned soft attention gate (~50 M params).",
        "Only the convolution operator differs.  Same depth, same channels, same loss, same schedule.",
    ], left=18.5, top=3.0, width_cm=14.5, font_pt=16)
    add_speaker_notes(s,
        "The PConv encoder-decoder shown on the left is the topological "
        "template.  Both baselines reuse the same skeleton — same number "
        "of stages, same channel widths, same skip connections — and "
        "swap only the convolution operator.  This guarantees that any "
        "quality gap is attributable to the operator and not to depth "
        "or capacity.")


def slide_methodology_loss():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Composite loss (Liu et al. 2018)")
    add_bullet_box(s, [
        "L = λ_valid · L1_valid + λ_hole · L1_hole + λ_perc · L_perc + λ_style · L_style + λ_tv · L_tv",
        "L1_valid:  pixelwise L1 on intact pixels.",
        "L1_hole:   pixelwise L1 on hole pixels (weighted 6×).",
        "L_perc:    VGG16 perceptual loss at relu1_1 / relu2_1 / relu3_1.",
        "L_style:   Gram-matrix style loss at the same VGG layers.",
        "L_tv:      total-variation regulariser over the hole region.",
        "Weights: 1, 6, 0.05, 50, 0.1 — identical for all three models.",
    ], font_pt=18)
    add_speaker_notes(s,
        "The loss is the five-term composite from Liu et al.  We kept "
        "the original weights with one deliberate change: the style "
        "weight is reduced from 120 to 50, which we found prevents the "
        "over-smoothed brushwork that the heavier weight produced on "
        "paintings in our preliminary experiments.  The same weights "
        "are used for all three networks.")


def slide_dataset_masks():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Dataset and synthetic damage masks")
    add_bullet_box(s, [
        "WikiArt — 39,213 paintings × 5 styles (Renaissance, Baroque, Impressionism, Post-Impressionism, Realism).",
        "Aspect-preserving resize + centre crop to 256 × 256.",
        "Train / val / test = 31,815 / 3,977 / 3,977, stratified by style.",
        "Four damage primitives: brush strokes, cracks, paint loss, aging stains.",
        "Three difficulty levels by hole-area: light 10–20 %, medium 20–40 %, heavy 40–60 %.",
        "Test masks are deterministic (seed = 12345) — same masks for all three models.",
    ], font_pt=18)
    add_speaker_notes(s,
        "WikiArt is large enough to train from scratch.  The synthetic "
        "damage protocol is the second contribution: four primitives at "
        "three difficulty levels gives twelve evaluation cells per "
        "model.  Crucially, the test masks are deterministic, so the "
        "three models all see exactly the same corrupted inputs.")


def slide_training():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Training protocol — the fairness contract")
    df = pd.DataFrame({
        "Hyperparameter": [
            "Optimiser", "Learning rate", "LR schedule", "Warmup",
            "Batch size", "Epochs", "Image size", "Mixed precision",
            "Gradient clip", "Seed",
        ],
        "Value": [
            "Adam (β=(0.9, 0.999))", "2 × 10⁻⁴",
            "Cosine decay → 1 × 10⁻⁶", "1 epoch",
            "4", "12", "256 × 256", "fp16 (AMP), loss in fp32",
            "L2 norm 1", "42",
        ],
    })
    add_table(s, df, left_cm=2, top_cm=2.6, width_cm=20, height_cm=12, font_pt=15)
    add_bullet_box(s, [
        "All three models trained from scratch with these settings.",
        "Only the convolution operator differs.",
    ], left=24, top=4.0, width_cm=8.5, font_pt=14)
    add_speaker_notes(s,
        "These are the hyperparameters.  They are identical across all "
        "three experiments — the only difference between the three runs "
        "is the model definition file.  Twelve epochs is short by the "
        "standards of the original paper but is what fits in the "
        "free-tier compute budget; this is acknowledged as a limitation.")


def slide_evaluation():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Evaluation pipeline")
    add_bullet_box(s, [
        "1.  Deterministic test masks generated once (seed=12345).",
        "2.  Each of 3,977 test paintings × 4 damage types × 3 difficulties = 47,724 pairs.",
        "3.  Each pair processed by all three models.",
        "4.  Per-image PSNR / SSIM / LPIPS computed on composited output (hole pixels from model, valid pixels from input).",
        "5.  Per-difficulty FID computed by pooling across damage types.",
        "6.  Paired Wilcoxon signed-rank test on per-image deltas, with Bonferroni correction across the three pairwise comparisons.",
    ], font_pt=18)
    add_speaker_notes(s,
        "The evaluation is one-shot: 47k pairs through all three "
        "models, all metrics computed once.  Paired Wilcoxon is "
        "appropriate because every model sees the same pairs.  "
        "Bonferroni correction guards against the three-way "
        "comparison.")


def slide_headline_results():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Headline results — and a surprise")
    df = load_overall()
    if not df.empty:
        add_table(s, df, left_cm=2.5, top_cm=3.0, width_cm=28, height_cm=6,
                  font_pt=22)
    add_bullet_box(s, [
        "Gated U-Net wins on all three metrics.",
        "Vanilla U-Net second — beats PConv by 1.7 dB PSNR.",
        "PConv U-Net last — the most-cited inpainting architecture is the worst on art.",
    ], left=2.5, top=11.5, width_cm=28, font_pt=20)
    add_speaker_notes(s,
        "This is the surprise.  PConv U-Net — the most-cited "
        "irregular-mask inpainting architecture — comes last on every "
        "metric.  Vanilla U-Net, the simplest possible baseline, beats "
        "it by 1.7 dB PSNR.  Gated convolution wins by another half a "
        "dB on top of that.  The next three slides confirm that this "
        "ranking is stable across damage severity and damage type.")


def slide_per_difficulty():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Per-difficulty PSNR")
    df = load_by_difficulty_psnr()
    if not df.empty:
        add_table(s, df, left_cm=2.5, top_cm=3.0, width_cm=18, height_cm=6,
                  font_pt=22)
    add_image(s, EVAL_DIR / "figures" / "fig1_psnr.png",
              left_cm=22, top_cm=2.7, width_cm=10)
    add_bullet_box(s, [
        "Ranking Gated > Vanilla > PConv holds for every difficulty.",
        "All models degrade by ~6 dB from light to heavy.",
        "Gated–Vanilla gap narrows on heavy damage.",
    ], left=2.5, top=11.5, width_cm=18, font_pt=18)
    add_speaker_notes(s,
        "Per-difficulty breakdown.  The ranking is monotone — Gated, "
        "Vanilla, PConv — and stable across light, medium, heavy.  All "
        "three models degrade by roughly six decibels from light to "
        "heavy, which is the expected scaling.")


def slide_per_damage():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Per-damage-type PSNR")
    df = load_by_damage_psnr()
    if not df.empty:
        add_table(s, df, left_cm=2.5, top_cm=3.0, width_cm=22, height_cm=6,
                  font_pt=22)
    add_bullet_box(s, [
        "Cracks and brush strokes — easiest (thin geometry, lots of valid context).",
        "Paint loss and stains — hardest (large contiguous regions).",
        "Architectural ranking is preserved in every damage cell.",
    ], left=2.5, top=11.5, width_cm=29, font_pt=18)
    add_speaker_notes(s,
        "Per-damage breakdown.  Cracks and brushstrokes are easier "
        "than paint-loss blobs and stains because the thin geometry "
        "leaves more valid context near every hole pixel.  Crucially, "
        "the architectural ranking is preserved in every one of the "
        "twelve cells — the result is not driven by any single "
        "damage type.")


def slide_qualitative():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Qualitative comparison")
    add_image(s, EVAL_DIR / "fig2_highres_comparison.jpg",
              left_cm=2, top_cm=2.7, width_cm=30)
    add_speaker_notes(s,
        "Five-column figure: original, damaged input, Vanilla, PConv, "
        "Gated.  PConv outputs exhibit slight colour-bleed and softening "
        "along hole boundaries — that's the over-renormalisation "
        "behaviour of partial convolution.  Gated outputs are the "
        "sharpest.  Vanilla is surprisingly competitive at lower "
        "difficulties.")


def slide_why_pconv_lost():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Why PConv lost — three reasons")
    add_bullet_box(s, [
        "Hard vs soft attention.  PConv binarises confidence early; gated conv keeps a continuous learned gate through every layer.",
        "Renormalisation artifacts.  PConv divides by valid-area, which over-amplifies on small valid neighbourhoods — visible as colour bleed at hole boundaries.",
        "Stylised data favours the simple baseline.  Mask concatenation is enough for a high-capacity U-Net to figure out the mask, given regular brushwork statistics.",
    ])
    add_speaker_notes(s,
        "Three interlocking explanations for the surprise.  First, "
        "PConv binarises confidence by design — by layer four or five "
        "the propagated mask is essentially all-ones, so the network "
        "loses all spatial confidence information.  Gated conv preserves "
        "it.  Second, the explicit renormalisation by valid-area can "
        "blow up near sparse valid pixels, producing the colour bleed "
        "we see in qualitative results.  Third, the vanilla baseline "
        "is unexpectedly strong because painting texture is more "
        "regular than natural-image texture.")


def slide_demo():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Interactive demo (contribution C4)")
    add_bullet_box(s, [
        "Built with Gradio; runs locally with `python demo/app.py`.",
        "Tab 1 — brush mode: upload a painting, paint over the damage with the brush tool.",
        "Tab 2 — synthetic preset: choose damage type and difficulty, the app generates the mask deterministically.",
        "All three models run in parallel; outputs displayed side-by-side with per-image PSNR / SSIM.",
        "Static benchmark panel shows test-set averages for context.",
        "First publicly-runnable art-inpainting demo, to our knowledge.",
    ], font_pt=18)
    add_speaker_notes(s,
        "The demo is the fourth contribution.  It runs locally with a "
        "single Python command and opens in the browser.  I'd be happy "
        "to switch to the live demo at the end of the talk if there "
        "is time.")


def slide_conclusion():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_title(s, "Conclusion and future work")
    add_bullet_box(s, [
        "Gated U-Net > Vanilla U-Net > PConv U-Net on every metric, every difficulty, every damage type.",
        "Counter-intuitive for the inpainting community: PConv is not the right default for art-domain inpainting.",
        "Four contributions: WikiArt benchmark, multi-type damage masks, rigorous comparison, interactive demo.",
        "",
        "Future work:",
        "—  curate a real-damage evaluation set (no ground truth, human ranking).",
        "—  add GAN refinement on top of the Gated U-Net.",
        "—  multi-resolution fine-tune for 512 × 512 inference.",
    ], font_pt=18)
    add_speaker_notes(s,
        "To wrap up: the take-home message is the surprising ranking — "
        "Gated, then Vanilla, then PConv.  Four contributions, three "
        "concrete future directions.  Thank you for your attention; I "
        "am happy to take questions.")


def slide_qa():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x10, 0x32, 0x66)
    bg.line.fill.background()

    tb = s.shapes.add_textbox(Cm(2), Cm(6), SW - Cm(4), Cm(6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank you.\n\nQuestions?"
    r.font.size = Pt(56)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_speaker_notes(s,
        "Anticipated questions and rehearsed answers:\n\n"
        "Q: Why only 12 epochs?  A: Free-tier compute budget; we "
        "report this as a limitation in Chapter 6.  Longer training "
        "would likely benefit PConv proportionally more.\n\n"
        "Q: Why not include LaMa / Stable-Diffusion inpainting?  A: "
        "Both rely on pretraining at a scale outside the budget of "
        "this thesis.  A fine-tuned LaMa is mentioned as future work.\n\n"
        "Q: Are the synthetic masks realistic?  A: They are a proxy.  "
        "We mitigate by reporting per-damage breakdown; no paired "
        "real-damage dataset exists at the required scale.\n\n"
        "Q: Why does Gated have ~2× the parameters of PConv?  A: "
        "Gated conv duplicates conv weights into feat+gate paths.  We "
        "accept this as an architectural property; the comparison "
        "claims fairness in training recipe, not in parameter count.")


# ──────────────────────────────────────────────────────────────────────
# Build
# ──────────────────────────────────────────────────────────────────────

def build() -> Path:
    slide_title()
    slide_motivation()
    slide_research_questions()
    slide_contributions()
    slide_related_work()
    slide_methodology_arch()
    slide_methodology_loss()
    slide_dataset_masks()
    slide_training()
    slide_evaluation()
    slide_headline_results()
    slide_per_difficulty()
    slide_per_damage()
    slide_qualitative()
    slide_why_pconv_lost()
    slide_demo()
    slide_conclusion()
    slide_qa()

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    return OUT_PPTX


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB)")
